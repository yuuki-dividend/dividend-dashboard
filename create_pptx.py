"""High Dividend Portfolio Tracker - PowerPoint Report Generator"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.chart.data import CategoryChartData
import os

OUTPUT = os.path.join(os.path.dirname(os.path.abspath(__file__)), "portfolio_report.pptx")

# Color palette
NAVY = RGBColor(0x1E, 0x27, 0x61)
DARK_BG = RGBColor(0x0F, 0x17, 0x2A)
ICE_BLUE = RGBColor(0xCA, 0xDC, 0xFC)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
ACCENT_BLUE = RGBColor(0x3B, 0x82, 0xF6)
ACCENT_GREEN = RGBColor(0x10, 0xB9, 0x81)
ACCENT_RED = RGBColor(0xEF, 0x44, 0x44)
ACCENT_GOLD = RGBColor(0xF5, 0x9E, 0x0B)
LIGHT_GRAY = RGBColor(0x94, 0xA3, 0xB8)
CARD_BG = RGBColor(0x1E, 0x29, 0x3B)
TEAL = RGBColor(0x06, 0xB6, 0xD4)
PURPLE = RGBColor(0x8B, 0x5C, 0xF6)


def add_rect(slide, left, top, width, height, fill_color):
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape


def set_bg(slide, color):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def text_box(slide, left, top, width, height, text, size=14,
             color=WHITE, bold=False, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = "Calibri"
    p.alignment = align
    return tb


def add_card(slide, x, y, w, h, title, value, sub="", val_color=WHITE, accent=ACCENT_BLUE):
    add_rect(slide, x, y, w, h, CARD_BG)
    add_rect(slide, x, y, w, Inches(0.04), accent)
    text_box(slide, x + Inches(0.15), y + Inches(0.35), w - Inches(0.3), Inches(0.2),
             title, size=9, color=LIGHT_GRAY)
    text_box(slide, x + Inches(0.15), y + Inches(0.55), w - Inches(0.3), Inches(0.35),
             value, size=18, color=val_color, bold=True)
    if sub:
        text_box(slide, x + Inches(0.15), y + Inches(0.95), w - Inches(0.3), Inches(0.2),
                 sub, size=8, color=LIGHT_GRAY)


def build():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # === SLIDE 1: Title ===
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, DARK_BG)
    add_rect(s, Inches(0), Inches(0), Inches(13.333), Inches(2.5), NAVY)
    add_rect(s, Inches(0), Inches(2.45), Inches(13.333), Inches(0.06), ACCENT_BLUE)

    text_box(s, Inches(1), Inches(0.8), Inches(11), Inches(0.8),
             "High Dividend Portfolio Tracker", size=42, bold=True, align=PP_ALIGN.CENTER)
    text_box(s, Inches(1), Inches(1.7), Inches(11), Inches(0.5),
             "53 Stocks | Annual Dividend: 85,302 JPY | Yield: 4.34%",
             size=16, color=ICE_BLUE, align=PP_ALIGN.CENTER)

    cards = [
        ("Total Cost", "1,964,891 JPY", "53 stocks", WHITE, ACCENT_BLUE),
        ("Market Value", "2,554,009 JPY", "Current eval", WHITE, ACCENT_BLUE),
        ("Gain/Loss", "+589,118 JPY", "+29.98%", ACCENT_GREEN, ACCENT_GREEN),
        ("Annual Div", "85,302 JPY", "After tax: 67,973", ACCENT_GOLD, ACCENT_GOLD),
        ("Avg Yield", "4.34%", "Current: 3.34%", ACCENT_BLUE, ACCENT_BLUE),
    ]
    for i, (t, v, sub, vc, ac) in enumerate(cards):
        add_card(s, Inches(0.8) + i * Inches(2.4), Inches(3.2),
                 Inches(2.2), Inches(1.3), t, v, sub, vc, ac)

    text_box(s, Inches(1), Inches(5.5), Inches(11), Inches(0.8),
             "All 53 stocks in NISA | Growth allowance 81.9% used (remaining: 435,109 JPY)",
             size=12, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

    # === SLIDE 2: Portfolio Overview ===
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, DARK_BG)
    add_rect(s, Inches(0), Inches(0), Inches(13.333), Inches(1.0), NAVY)
    text_box(s, Inches(0.8), Inches(0.25), Inches(8), Inches(0.5),
             "Portfolio Overview", size=28, bold=True)

    metrics = [
        ("Total Stocks", "53", ACCENT_BLUE),
        ("Total Cost", "1,964,891 JPY", ACCENT_BLUE),
        ("Market Value", "2,554,009 JPY", ACCENT_BLUE),
        ("Unrealized Gain", "+589,118 JPY (+29.98%)", ACCENT_GREEN),
        ("Annual Div (Pre-tax)", "85,302 JPY", ACCENT_GOLD),
        ("Annual Div (After-tax)", "67,973 JPY", ACCENT_GOLD),
        ("Avg Cost Yield", "4.34%", ACCENT_GREEN),
        ("Avg Current Yield", "3.34%", ACCENT_GREEN),
    ]
    for i, (label, val, clr) in enumerate(metrics):
        y = Inches(1.3) + i * Inches(0.72)
        add_rect(s, Inches(0.5), y, Inches(6), Inches(0.62), CARD_BG)
        add_rect(s, Inches(0.5), y, Inches(0.06), Inches(0.62), clr)
        text_box(s, Inches(0.7), y + Inches(0.15), Inches(2.8), Inches(0.3),
                 label, size=11, color=LIGHT_GRAY)
        vc = ACCENT_GREEN if "+" in val else WHITE
        text_box(s, Inches(3.5), y + Inches(0.15), Inches(2.8), Inches(0.3),
                 val, size=13, bold=True, color=vc, align=PP_ALIGN.RIGHT)

    # NISA card
    nx = Inches(7)
    add_rect(s, nx, Inches(1.3), Inches(5.8), Inches(2.8), CARD_BG)
    add_rect(s, nx, Inches(1.3), Inches(5.8), Inches(0.05), ACCENT_GOLD)
    text_box(s, nx + Inches(0.3), Inches(1.45), Inches(5), Inches(0.35),
             "NISA Account Status", size=16, bold=True, color=ACCENT_GOLD)
    nisa = [("Growth Used", "81.9%"), ("Remaining", "435,109 JPY"),
            ("Lifetime Used", "16.4%"), ("Tax Benefit", "NISA tax-free")]
    for i, (l, v) in enumerate(nisa):
        ny = Inches(2.0) + i * Inches(0.5)
        text_box(s, nx + Inches(0.3), ny, Inches(2.8), Inches(0.3), l, size=10, color=LIGHT_GRAY)
        text_box(s, nx + Inches(3.3), ny, Inches(2.2), Inches(0.3),
                 v, size=12, bold=True, color=ACCENT_BLUE, align=PP_ALIGN.RIGHT)

    # Performance
    add_rect(s, nx, Inches(4.5), Inches(5.8), Inches(2.3), CARD_BG)
    add_rect(s, nx, Inches(4.5), Inches(5.8), Inches(0.05), ACCENT_GREEN)
    text_box(s, nx + Inches(0.3), Inches(4.65), Inches(5), Inches(0.3),
             "Performance", size=16, bold=True, color=ACCENT_GREEN)
    text_box(s, nx + Inches(0.3), Inches(5.1), Inches(5), Inches(0.7),
             "+29.98%", size=48, bold=True, color=ACCENT_GREEN, align=PP_ALIGN.CENTER)
    text_box(s, nx + Inches(0.3), Inches(5.9), Inches(5), Inches(0.3),
             "+589,118 JPY unrealized gain", size=12, align=PP_ALIGN.CENTER)

    # === SLIDE 3: Monthly Dividends ===
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, DARK_BG)
    add_rect(s, Inches(0), Inches(0), Inches(13.333), Inches(1.0), NAVY)
    text_box(s, Inches(0.8), Inches(0.25), Inches(8), Inches(0.5),
             "Monthly Dividend Calendar", size=28, bold=True)

    cd = CategoryChartData()
    cd.categories = ['Jan', 'Feb', 'Mar', 'Apr', 'May', 'Jun',
                     'Jul', 'Aug', 'Sep', 'Oct', 'Nov', 'Dec']
    vals = [3200, 0, 8500, 2800, 0, 39978, 3200, 0, 8500, 2800, 0, 35890]
    cd.add_series('Pre-tax', vals)
    cd.add_series('After-tax', [int(v * 0.797) for v in vals])

    cf = s.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED,
                            Inches(0.5), Inches(1.2), Inches(8), Inches(5.5), cd)
    ch = cf.chart
    ch.has_legend = True
    ch.legend.position = XL_LEGEND_POSITION.BOTTOM
    ch.legend.font.size = Pt(9)
    ch.legend.font.color.rgb = LIGHT_GRAY
    ch.plots[0].gap_width = 80
    ch.plots[0].series[0].format.fill.solid()
    ch.plots[0].series[0].format.fill.fore_color.rgb = ACCENT_BLUE
    ch.plots[0].series[1].format.fill.solid()
    ch.plots[0].series[1].format.fill.fore_color.rgb = TEAL

    sx = Inches(9)
    for i, (title, val, clr) in enumerate([
        ("Peak: June", "39,978 JPY", ACCENT_GOLD),
        ("Peak: December", "35,890 JPY", ACCENT_GOLD),
        ("Monthly Average", "7,109 JPY", ACCENT_BLUE),
        ("Annual Total", "85,302 JPY", ACCENT_GREEN),
    ]):
        cy = Inches(1.3) + i * Inches(1.4)
        add_rect(s, sx, cy, Inches(3.8), Inches(1.2), CARD_BG)
        add_rect(s, sx, cy, Inches(3.8), Inches(0.04), clr)
        text_box(s, sx + Inches(0.2), cy + Inches(0.12), Inches(3.4), Inches(0.3),
                 title, size=12, bold=True, color=clr)
        text_box(s, sx + Inches(0.2), cy + Inches(0.5), Inches(3.4), Inches(0.4),
                 val, size=20, bold=True)

    # === SLIDE 4: Sector Composition ===
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, DARK_BG)
    add_rect(s, Inches(0), Inches(0), Inches(13.333), Inches(1.0), NAVY)
    text_box(s, Inches(0.8), Inches(0.25), Inches(8), Inches(0.5),
             "Sector Composition", size=28, bold=True)

    sd = CategoryChartData()
    sd.categories = ['Wholesale', 'Machinery', 'IT/Telecom', 'Banking',
                     'Chemicals', 'Construction', 'Real Estate', 'Others']
    sd.add_series('Weight', [14.3, 10.6, 8.5, 7.8, 7.2, 6.5, 5.8, 39.3])
    cf = s.shapes.add_chart(XL_CHART_TYPE.PIE, Inches(0.5), Inches(1.2),
                            Inches(5.5), Inches(5.5), sd)
    ch = cf.chart
    ch.has_legend = True
    ch.legend.position = XL_LEGEND_POSITION.BOTTOM
    ch.legend.font.size = Pt(9)
    ch.legend.font.color.rgb = LIGHT_GRAY
    colors = ["3B82F6", "10B981", "F59E0B", "8B5CF6", "EF4444", "06B6D4", "EC4899", "64748B"]
    for i, c in enumerate(colors):
        pt = ch.plots[0].series[0].points[i]
        pt.format.fill.solid()
        pt.format.fill.fore_color.rgb = RGBColor(int(c[:2], 16), int(c[2:4], 16), int(c[4:], 16))

    rx = Inches(6.5)
    add_rect(s, rx, Inches(1.3), Inches(6.3), Inches(2.2), CARD_BG)
    add_rect(s, rx, Inches(1.3), Inches(6.3), Inches(0.04), PURPLE)
    text_box(s, rx + Inches(0.3), Inches(1.4), Inches(5), Inches(0.3),
             "Defensive vs Cyclical", size=16, bold=True)

    bx = rx + Inches(0.3)
    text_box(s, bx, Inches(1.9), Inches(2), Inches(0.25), "Defensive 27.4%", size=11, color=TEAL, bold=True)
    add_rect(s, bx, Inches(2.2), Inches(5.5), Inches(0.3), CARD_BG)
    add_rect(s, bx, Inches(2.2), Inches(5.5 * 0.274), Inches(0.3), TEAL)
    text_box(s, bx, Inches(2.65), Inches(2), Inches(0.25), "Cyclical 72.6%", size=11, color=ACCENT_BLUE, bold=True)
    add_rect(s, bx, Inches(2.95), Inches(5.5), Inches(0.3), CARD_BG)
    add_rect(s, bx, Inches(2.95), Inches(5.5 * 0.726), Inches(0.3), ACCENT_BLUE)

    # Top sectors
    add_rect(s, rx, Inches(3.9), Inches(6.3), Inches(3.3), CARD_BG)
    add_rect(s, rx, Inches(3.9), Inches(6.3), Inches(0.04), ACCENT_GOLD)
    text_box(s, rx + Inches(0.3), Inches(4.0), Inches(5), Inches(0.3),
             "Top Sectors", size=14, bold=True, color=ACCENT_GOLD)
    for i, (name, pct) in enumerate([
        ("Wholesale (Shosha)", "14.3%"), ("Machinery", "10.6%"),
        ("IT & Telecom", "8.5%"), ("Banking", "7.8%"),
        ("Chemicals", "7.2%"), ("Construction", "6.5%"), ("Real Estate", "5.8%")
    ]):
        sy = Inches(4.45) + i * Inches(0.37)
        text_box(s, rx + Inches(0.3), sy, Inches(3.5), Inches(0.3), name, size=10)
        text_box(s, rx + Inches(4), sy, Inches(2), Inches(0.3),
                 pct, size=11, bold=True, color=ACCENT_BLUE, align=PP_ALIGN.RIGHT)

    # === SLIDE 5: Dividend Scorecard ===
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, DARK_BG)
    add_rect(s, Inches(0), Inches(0), Inches(13.333), Inches(1.0), NAVY)
    text_box(s, Inches(0.8), Inches(0.25), Inches(8), Inches(0.5),
             "Dividend Scorecard", size=28, bold=True)

    # Average card
    add_rect(s, Inches(0.5), Inches(1.3), Inches(3.5), Inches(2.5), CARD_BG)
    add_rect(s, Inches(0.5), Inches(1.3), Inches(3.5), Inches(0.04), ACCENT_GOLD)
    text_box(s, Inches(0.8), Inches(1.4), Inches(3), Inches(0.3),
             "Portfolio Average", size=14, bold=True, color=ACCENT_GOLD)
    text_box(s, Inches(0.8), Inches(1.85), Inches(3), Inches(0.7),
             "C", size=56, bold=True, color=ACCENT_GOLD, align=PP_ALIGN.CENTER)
    text_box(s, Inches(0.8), Inches(2.7), Inches(3), Inches(0.3),
             "Score: 49 / 100", size=14, align=PP_ALIGN.CENTER)

    grades = [("S", "2", ACCENT_GREEN), ("A", "8", ACCENT_BLUE), ("B", "15", TEAL),
              ("C", "18", ACCENT_GOLD), ("D", "10", ACCENT_RED)]
    for i, (g, cnt, clr) in enumerate(grades):
        gx = Inches(4.5) + i * Inches(1.7)
        add_rect(s, gx, Inches(1.3), Inches(1.5), Inches(2.5), CARD_BG)
        add_rect(s, gx, Inches(1.3), Inches(1.5), Inches(0.04), clr)
        text_box(s, gx, Inches(1.5), Inches(1.5), Inches(0.6),
                 g, size=36, bold=True, color=clr, align=PP_ALIGN.CENTER)
        text_box(s, gx, Inches(2.2), Inches(1.5), Inches(0.3),
                 cnt + " stocks", size=11, bold=True, align=PP_ALIGN.CENTER)

    # Top stocks
    add_rect(s, Inches(0.5), Inches(4.2), Inches(12.3), Inches(3.0), CARD_BG)
    add_rect(s, Inches(0.5), Inches(4.2), Inches(12.3), Inches(0.04), ACCENT_GREEN)
    text_box(s, Inches(0.8), Inches(4.3), Inches(5), Inches(0.3),
             "Top Rated Stocks", size=16, bold=True, color=ACCENT_GREEN)

    col_x = [0.8, 1.8, 5.5, 7.0, 8.5, 10.0]
    for i, h in enumerate(["#", "Stock", "Grade", "Score", "Yield", "Stability"]):
        text_box(s, Inches(col_x[i]), Inches(4.75), Inches(1.5), Inches(0.25),
                 h, size=9, bold=True, color=LIGHT_GRAY)

    rows = [
        ("1", "Okinawa Electric (9511)", "S", "89", "5.12%", "High"),
        ("2", "Mitsubishi HC Cap (8593)", "S", "83", "4.85%", "Very High"),
        ("3", "SMFG (8316)", "A", "78", "4.21%", "High"),
        ("4", "INPEX (1605)", "A", "75", "4.56%", "Medium"),
        ("5", "Japan Tobacco (2914)", "A", "73", "4.82%", "Very High"),
    ]
    for ri, row in enumerate(rows):
        ry = Inches(5.1) + ri * Inches(0.38)
        for ci, val in enumerate(row):
            clr = ACCENT_GREEN if ci == 2 and val == "S" else (ACCENT_BLUE if ci == 2 else WHITE)
            text_box(s, Inches(col_x[ci]), ry, Inches(1.5 if ci != 1 else 3.5), Inches(0.25),
                     val, size=10, color=clr, bold=(ci in [2, 3]))

    # === SLIDE 6: App Features ===
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, DARK_BG)
    add_rect(s, Inches(0), Inches(0), Inches(13.333), Inches(1.0), NAVY)
    text_box(s, Inches(0.8), Inches(0.25), Inches(8), Inches(0.5),
             "App Features", size=28, bold=True)

    features = [
        ("Dashboard", "Real-time portfolio summary with charts", ACCENT_BLUE),
        ("Portfolio Search", "Sort, filter by sector, yield, score", ACCENT_GREEN),
        ("Dividend Calendar", "12-month view with timeline", TEAL),
        ("Buy Timing Tool", "Screening for undervalued stocks", ACCENT_GOLD),
        ("Watchlist", "Track stocks before buying", PURPLE),
        ("Stock Checker", "Deep analysis and scoring", ACCENT_BLUE),
        ("NISA Manager", "Track allowance and tax savings", ACCENT_GREEN),
        ("DRIP Simulator", "Dividend reinvestment projection", TEAL),
        ("P&L Heatmap", "Visual gain/loss overview", ACCENT_GOLD),
        ("Yield Alerts", "Notifications for yield targets", ACCENT_RED),
        ("Export", "CSV, JSON, Google Sheets sync", PURPLE),
    ]
    cw, ch_h = Inches(3.8), Inches(1.2)
    for i, (name, desc, clr) in enumerate(features):
        col, row = i % 3, i // 3
        fx = Inches(0.5) + col * Inches(4.05)
        fy = Inches(1.3) + row * Inches(1.4)
        add_rect(s, fx, fy, cw, ch_h, CARD_BG)
        add_rect(s, fx, fy, Inches(0.06), ch_h, clr)
        text_box(s, fx + Inches(0.25), fy + Inches(0.15), cw - Inches(0.4), Inches(0.3),
                 name, size=13, bold=True, color=clr)
        text_box(s, fx + Inches(0.25), fy + Inches(0.55), cw - Inches(0.4), Inches(0.5),
                 desc, size=9, color=LIGHT_GRAY)

    # === SLIDE 7: Future Improvements ===
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, DARK_BG)
    add_rect(s, Inches(0), Inches(0), Inches(13.333), Inches(1.0), NAVY)
    text_box(s, Inches(0.8), Inches(0.25), Inches(8), Inches(0.5),
             "Future Improvements", size=28, bold=True)

    improvements = [
        ("Real-time Price Updates", "Auto stock price refresh via API (stooq, Yahoo Finance)", ACCENT_BLUE),
        ("Sector Rebalancing", "AI suggestions for portfolio rebalancing", ACCENT_GREEN),
        ("Dividend Growth Tracking", "Historical growth rate and future projections", TEAL),
        ("Mobile Responsive", "Full mobile optimization for on-the-go monitoring", PURPLE),
        ("Multi-currency", "US stocks and ADRs with FX conversion", ACCENT_GOLD),
        ("Tax Optimization", "Optimal NISA allocation and tax-loss harvesting", ACCENT_RED),
    ]
    for i, (title, desc, clr) in enumerate(improvements):
        col, row = i % 2, i // 2
        ix = Inches(0.5) + col * Inches(6.4)
        iy = Inches(1.3) + row * Inches(1.9)
        add_rect(s, ix, iy, Inches(6.1), Inches(1.7), CARD_BG)
        add_rect(s, ix, iy, Inches(6.1), Inches(0.04), clr)
        add_rect(s, ix + Inches(0.2), iy + Inches(0.2), Inches(0.4), Inches(0.4), clr)
        text_box(s, ix + Inches(0.2), iy + Inches(0.22), Inches(0.4), Inches(0.35),
                 str(i + 1), size=16, bold=True, align=PP_ALIGN.CENTER)
        text_box(s, ix + Inches(0.75), iy + Inches(0.2), Inches(5), Inches(0.3),
                 title, size=14, bold=True, color=clr)
        text_box(s, ix + Inches(0.75), iy + Inches(0.6), Inches(5), Inches(0.8),
                 desc, size=10, color=LIGHT_GRAY)

    # === SLIDE 8: Summary ===
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, DARK_BG)
    add_rect(s, Inches(0), Inches(0), Inches(13.333), Inches(3.0), NAVY)
    add_rect(s, Inches(0), Inches(2.95), Inches(13.333), Inches(0.06), ACCENT_BLUE)

    text_box(s, Inches(1), Inches(1.0), Inches(11), Inches(0.8),
             "High Dividend Portfolio Tracker", size=36, bold=True, align=PP_ALIGN.CENTER)
    text_box(s, Inches(1), Inches(1.8), Inches(11), Inches(0.5),
             "Building wealth through dividend income",
             size=16, color=ICE_BLUE, align=PP_ALIGN.CENTER)

    for i, (val, label) in enumerate([
        ("53", "Stocks"), ("4.34%", "Yield"), ("85,302", "Annual Div"),
        ("+29.98%", "Return"), ("11", "Features")
    ]):
        sx = Inches(1.5) + i * Inches(2.2)
        vc = ACCENT_GREEN if "+" in val else ACCENT_BLUE
        text_box(s, sx, Inches(4.0), Inches(2), Inches(0.6),
                 val, size=32, bold=True, color=vc, align=PP_ALIGN.CENTER)
        text_box(s, sx, Inches(4.6), Inches(2), Inches(0.3),
                 label, size=11, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

    text_box(s, Inches(1), Inches(5.8), Inches(11), Inches(0.4),
             "localhost:8080 | Python + Chart.js | Google Sheets Integration",
             size=12, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

    prs.save(OUTPUT)
    print(f"Saved: {OUTPUT}")


if __name__ == "__main__":
    build()
